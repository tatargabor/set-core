#!/usr/bin/env python3
"""
Build professional PPTX presentation for SET using python-pptx.
Uses Mermaid-rendered diagrams and dashboard screenshots.

Usage:
    python3 scripts/build-presentation.py [--lang en|hu] [--output FILE]
"""
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
ROOT = Path(__file__).resolve().parent.parent
PRES_DIR = ROOT / "docs" / "presentation"
DIAG_DIR = PRES_DIR / "diagrams"
IMG_DIR = ROOT / "docs" / "images" / "auto"
WEB_IMG = IMG_DIR / "web"
APP_IMG = IMG_DIR / "app"
CLI_IMG = IMG_DIR / "cli"

# ── Theme ──
BG = RGBColor(0x0F, 0x17, 0x2A)
FG = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x38, 0xBD, 0xF8)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
MUTED = RGBColor(0x6B, 0x72, 0x80)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
BORDER = RGBColor(0x33, 0x41, 0x55)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFA, 0xFA, 0xFA)

FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=FG, bold=False, align=PP_ALIGN.LEFT, font_name=FONT_BODY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=FG, font_name=FONT_BODY, line_spacing=1.3):
    """Add multiple paragraphs. lines = list of (text, color, bold)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text, c, bold = item if len(item) == 3 else (item[0], color, False)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = c
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox


def add_image(slide, path, left, top, width=None, height=None):
    p = Path(path)
    if not p.exists():
        # Add placeholder text instead
        add_text(slide, left, top, 4, 0.5, f"[missing: {p.name}]", 12, MUTED)
        return
    kwargs = {}
    if width: kwargs["width"] = Inches(width)
    if height: kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(p), Inches(left), Inches(top), **kwargs)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                        Inches(left), Inches(top),
                                        Inches(width), Inches(height))
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = BLUE
            p.font.bold = True
            p.font.name = FONT_BODY

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = FG
                p.font.name = FONT_BODY

    return tbl_shape


def title_slide(prs, title, subtitle, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    add_text(slide, 0, 1.5, 13.333, 1.5, title, 60, GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, 0, 3.2, 13.333, 1, subtitle, 24, AMBER, True, PP_ALIGN.CENTER)
    if note:
        add_text(slide, 0, 4.5, 13.333, 0.5, note, 16, MUTED, False, PP_ALIGN.CENTER)
    return slide


def section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(0x12, 0x1D, 0x33))
    add_text(slide, 0, 2.5, 13.333, 1.2, title, 44, BLUE, True, PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, 0, 4.0, 13.333, 0.8, subtitle, 20, MUTED, False, PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title, bullets=None, image_path=None, image_right=False,
                  note="", two_images=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.6, 0.3, 12, 0.7, title, 28, BLUE, True)

    if two_images:
        # Two images side by side
        left_img, right_img = two_images
        add_image(slide, left_img, 0.6, 1.5, width=5.8)
        add_image(slide, right_img, 6.8, 1.5, width=5.8)
    elif image_path and bullets:
        if image_right:
            add_multiline(slide, 0.6, 1.3, 5.5, 5, bullets, 14)
            add_image(slide, image_path, 6.8, 1.3, width=6)
        else:
            add_image(slide, image_path, 0.6, 1.3, width=6)
            add_multiline(slide, 7.0, 1.3, 5.5, 5, bullets, 14)
    elif image_path:
        add_image(slide, image_path, 0.8, 1.3, width=11.5)
    elif bullets:
        add_multiline(slide, 0.6, 1.3, 12, 5.5, bullets, 16)

    if note:
        add_text(slide, 0.6, 6.5, 12, 0.6, note, 13, AMBER, True)
    return slide


def diagram_slide(prs, title, diagram_name, note="", subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.6, 0.3, 12, 0.7, title, 28, BLUE, True)
    if subtitle:
        add_text(slide, 0.6, 0.9, 12, 0.5, subtitle, 16, MUTED)

    diag_path = DIAG_DIR / f"{diagram_name}.png"
    if diag_path.exists():
        add_image(slide, diag_path, 0.8, 1.6, width=11.5)
    else:
        add_text(slide, 2, 3, 8, 1, f"[diagram not found: {diagram_name}]", 18, RED)

    if note:
        add_text(slide, 0.6, 6.5, 12, 0.6, note, 13, AMBER, True)
    return slide


def table_slide(prs, title, headers, rows, col_widths=None, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.6, 0.3, 12, 0.7, title, 28, BLUE, True)

    tw = sum(col_widths) if col_widths else 11
    add_table(slide, 0.6, 1.3, tw, 0.4 * (len(rows) + 1), headers, rows, col_widths)

    if note:
        add_text(slide, 0.6, 6.5, 12, 0.6, note, 13, AMBER, True)
    return slide


def build_en(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── 1. Title ──
    title_slide(prs, "SET", "Autonomous multi-agent orchestration for Claude Code",
                "Give it a spec — get merged features.  |  April 2026")

    # ── 2. Agenda ──
    content_slide(prs, "Agenda", [
        ("The problem — why prompt and pray isn't enough", MUTED, False),
        ("Six pillars — SPECIFY, DECOMPOSE, EXECUTE, SUPERVISE, VERIFY, LEARN", GREEN, True),
        ("The input — spec + design = quality input", MUTED, False),
        ("E2E Demo: MiniShop — the full pipeline walkthrough", MUTED, False),
        ("Works everywhere — greenfield, brownfield, isolated unit", MUTED, False),
        ("Dashboard & monitoring — real-time visibility", MUTED, False),
        ("Lessons learned — 100+ runs in production", MUTED, False),
        ("Roadmap — where we're heading", MUTED, False),
    ])

    # ── 3. The problem ──
    section_slide(prs, "The Problem", "Why isn't prompting enough?")

    table_slide(prs, "AI coding today — the challenges",
        ["Problem", "Typical approach", "Result"],
        [
            ["Divergence", "Run prompt twice = 2 different results", "Not reproducible"],
            ["Hallucination", "Agents make up what they don't know", "Missing features"],
            ["Quality roulette", "LLM judges code quality", "Inconsistent"],
            ["Spec drift", "\"Tests pass\" ≠ \"spec is satisfied\"", "Partial implementation"],
            ["Amnesia", "Every session starts from scratch", "Repeated mistakes"],
            ["Bug fixing", "Manual debugging, hours", "Lost time"],
        ], col_widths=[2.5, 4.5, 4],
        note="Most AI coding tools are non-deterministic — same prompt, different result.")

    table_slide(prs, "How SET addresses it",
        ["Challenge", "SET solution", "Result"],
        [
            ["Divergence", "3-layer template system", "83-87% convergence"],
            ["Hallucination", "OpenSpec + acceptance criteria", "Implements against spec"],
            ["Quality roulette", "Programmatic gates (exit codes)", "Deterministic pass/fail"],
            ["Spec drift", "Coverage tracking + auto-replan", "100% spec coverage"],
            ["Amnesia", "Hook-driven memory (5 layers)", "100% context capture"],
            ["Bug fixing", "Issue pipeline: detect → fix", "30 second recovery"],
        ], col_widths=[2.5, 4.5, 4],
        note="\"We don't prompt — we specify.\"")

    # ── 4. Six pillars ──
    section_slide(prs, "Six Pillars", "The architecture behind autonomous orchestration")

    diagram_slide(prs, "SPECIFY → DECOMPOSE → EXECUTE → SUPERVISE → VERIFY → LEARN",
                  "01-six-pillars",
                  "Every feature, gate, and automation maps to one of these six pillars.")

    diagram_slide(prs, "Change lifecycle — state machine", "02-state-machine",
                  "The system doesn't stop until the spec is 100% covered.")

    diagram_slide(prs, "Three-tier supervision", "08-supervision",
                  "Agents handle code errors. Orchestrator handles workflow. Sentinel handles infrastructure.")

    diagram_slide(prs, "Gradual escalation — when things break", "05-escalation",
                  "L2 (restart) solves 70% of problems. L4 (give up) is better than burning unlimited tokens.")

    # ── 5. The input ──
    section_slide(prs, "The Input", "Output quality depends on input quality")

    content_slide(prs, "What a good spec looks like", [
        ("Data model — entities, fields, relations, enums", FG, False),
        ("Pages — sections, columns, components", FG, False),
        ("Design tokens — hex colors, fonts, spacing", FG, False),
        ("Auth & roles — protected routes, registration", FG, False),
        ("Seed data — realistic names, not \"Product 1\"", FG, False),
        ("i18n — locales, URL structure", FG, False),
        ("Business requirements — user stories, acceptance criteria", FG, False),
        ("", FG, False),
        ("You are the product owner. Agents are the dev team.", AMBER, True),
        ("The spec is the sprint backlog.", AMBER, True),
    ])

    content_slide(prs, "Spec + Figma Design",
                  two_images=(CLI_IMG / "spec-preview.png",
                              Path(ROOT / "docs/images/auto/figma/storefront-design.png")),
                  note="set-design-sync extracts Figma tokens → design-system.md → agents read before implementation")

    # ── 6. E2E Demo ──
    section_slide(prs, "E2E Demo: MiniShop", "Building a webshop from a spec — live")

    diagram_slide(prs, "The pipeline", "03-pipeline",
                  "7 stages. Fully automatic. Sentinel supervises everything.")

    diagram_slide(prs, "Dependency DAG — phased execution", "04-dag",
                  "Phase 2 runs in parallel: products + admin-auth have no dependencies on each other.")

    content_slide(prs, "Step 1: Digest — extracting requirements",
                  image_path=WEB_IMG / "tab-digest.png",
                  bullets=[
                      ("Domains: Products, Cart, Auth, Admin", FG, False),
                      ("32 requirements with REQ-IDs", FG, False),
                      ("84 acceptance criteria (WHEN/THEN)", FG, False),
                      ("", FG, False),
                      ("Every REQ-ID tracked through", GREEN, True),
                      ("the entire pipeline.", GREEN, True),
                  ], image_right=False)

    content_slide(prs, "Step 3: Dispatch — parallel agents in worktrees",
                  bullets=[
                      ("Each change → own git worktree", FG, False),
                      ("Full isolation: can't interfere", FG, False),
                      ("Even with 1 change: worktree keeps main clean", FG, False),
                      ("", FG, False),
                      ("Inside each agent the Ralph Loop runs:", GREEN, True),
                      ("proposal → design → spec → tasks → code → test", FG, False),
                      ("", FG, False),
                      ("One agent = one change = one worktree", AMBER, True),
                  ],
                  image_path=WEB_IMG / "tab-sessions.png",
                  image_right=True)

    content_slide(prs, "The agent works — real-time terminal",
                  image_path=WEB_IMG / "tab-agent.png",
                  note="Not a black box — every step is observable: code writing, test running, bug fixing.")

    # ── Quality gates ──
    table_slide(prs, "Step 5: Verify — 7 quality gates",
        ["Gate", "Time", "Checks", "Type"],
        [
            ["Test", "8s", "Unit/integration tests", "Deterministic (exit code)"],
            ["Build", "35s", "Type check + bundle", "Deterministic (exit code)"],
            ["E2E", "45s", "Browser-based tests", "Deterministic (exit code)"],
            ["Review", "25s", "Code quality", "LLM (CRITICAL=fail)"],
            ["Coverage", "—", "Requirement coverage", "LLM + pattern match"],
            ["Smoke", "15s", "Post-merge sanity", "Deterministic"],
        ], col_widths=[2, 1.5, 4, 3.5],
        note="Total gate time: 422 seconds (12% of build time). Fast gates first — Jest fails? No waiting for Playwright.")

    diagram_slide(prs, "Self-healing pipeline", "07-self-healing",
                  "MiniShop: 5 gate failures, 5 autonomous fixes. Including IDOR vulnerabilities caught and patched.")

    table_slide(prs, "5 gate failures — 5 automatic fixes",
        ["#", "Failure", "Gate", "Agent fix"],
        [
            ["1", "Missing test file", "Test", "Added test file"],
            ["2", "Jest config error", "Build", "Fixed path mapping"],
            ["3", "Playwright auth tests (3 specs)", "E2E", "Updated redirects"],
            ["4", "Post-merge type error", "Build", "Synced with main"],
            ["5", "Cart test race condition", "E2E", "Added waitForSelector"],
        ], col_widths=[0.5, 3.5, 1.5, 5.5],
        note="Without gates, these 5 bugs would have merged into main and caused cascading failures.")

    # ── Merge + replan ──
    content_slide(prs, "Step 6: Merge — integration",
                  image_path=WEB_IMG / "tab-phases.png",
                  bullets=[
                      ("FF-only merge — clean history", FG, False),
                      ("Sequential merge queue", FG, False),
                      ("Phase ordering — phase 2 after 1", FG, False),
                      ("Post-merge smoke test", FG, False),
                      ("", FG, False),
                      ("The merge queue is the bottleneck", AMBER, True),
                      ("— intentionally.", AMBER, True),
                  ], image_right=True)

    # ── Works everywhere ──
    section_slide(prs, "Works Everywhere", "Not just \"build me an app from scratch\"")

    table_slide(prs, "Greenfield, Brownfield, Isolated Unit",
        ["Mode", "What it means", "Example"],
        [
            ["Greenfield", "Full app from spec + design", "MiniShop: 6/6, 1h 45m, 0 interventions"],
            ["Brownfield", "Existing codebase + new features", "SET builds itself: 1,500+ commits"],
            ["Isolated unit", "One module, one feature, one fix", "3 API endpoints with auth in existing Next.js"],
        ], col_widths=[2, 4.5, 4.5],
        note="Same pipeline in all modes. The entry barrier is not a 30-page spec — it can be a single task description.")

    # ── Results ──
    section_slide(prs, "The Result", "What the pipeline produced")

    content_slide(prs, "MiniShop — the finished application",
                  two_images=(APP_IMG / "products.png", APP_IMG / "product-detail.png"),
                  note="Product listing + detail pages — real data, working navigation, responsive layout")

    content_slide(prs, "MiniShop — cart and admin",
                  two_images=(APP_IMG / "cart.png", APP_IMG / "admin-dashboard.png"),
                  note="Working cart with totals | Admin dashboard — protected route, session-based auth")

    table_slide(prs, "The numbers — MiniShop benchmark",
        ["Metric", "Value"],
        [
            ["Changes planned/merged", "6/6 (100%)"],
            ["Total wall time", "1h 45m"],
            ["Human interventions", "0"],
            ["Merge conflicts", "0"],
            ["Jest unit tests", "38 (6 suites)"],
            ["Playwright E2E tests", "32 (6 spec files)"],
            ["Total tokens", "2.7M"],
            ["Gate retries", "5 (all self-healed)"],
        ], col_widths=[5.5, 5.5],
        note="Roughly equivalent to a day's work by 3-4 senior developers.")

    content_slide(prs, "Token usage",
                  image_path=WEB_IMG / "tab-tokens.png",
                  note="Cache ratio 26:1 — prompt caching dramatically reduces cost. Only input + output billed.")

    # ── Dashboard ──
    section_slide(prs, "Dashboard & Monitoring", "Real-time visibility into everything")

    content_slide(prs, "Web Dashboard — localhost:7400",
                  image_path=WEB_IMG / "dashboard-overview.png",
                  note="Changes, Phases, Tokens, Sessions, Sentinel, Log, Agent, Learnings, Digest — all live.")

    content_slide(prs, "Dashboard tabs — live data",
                  two_images=(WEB_IMG / "tab-changes.png", WEB_IMG / "tab-learnings.png"),
                  note="Changes: status at a glance | Learnings: gate failures, agent solutions, patterns")

    # ── Architecture ──
    section_slide(prs, "Architecture", "3 layers, extensible plugin system")

    diagram_slide(prs, "3-layer architecture", "06-architecture",
                  "Layer 1 never contains project-specific logic. All web code lives in Layer 2.")

    # ── Lessons ──
    section_slide(prs, "Lessons Learned", "100+ runs, real production experience")

    content_slide(prs, "8 lessons from production", [
        ("1. Agents need structure, not just prompts", GREEN, True),
        ("   OpenSpec artifacts keep them focused. Without: 3 agents, 3 table libs.", FG, False),
        ("2. Quality gates must be deterministic", GREEN, True),
        ("   Exit code > LLM judgment. Agents gamed the LLM review.", FG, False),
        ("3. Merge conflicts are #1 cascading failure", GREEN, True),
        ("   Phase ordering + DAG + sequential merge queue.", FG, False),
        ("4. Memory without hooks is useless", GREEN, True),
        ("   15+ sessions, 0 voluntary saves. 5-layer hooks: +34% improvement.", FG, False),
    ])

    content_slide(prs, "8 lessons (continued)", [
        ("5. E2E testing reveals what unit tests don't", GREEN, True),
        ("   Stale locks, zombie worktrees, race conditions in poll cycle.", FG, False),
        ("6. Stall detection needs grace periods", GREEN, True),
        ("   pnpm install >60s = watchdog killed it. Context-aware timeouts.", FG, False),
        ("7. The Sentinel pays for itself", GREEN, True),
        ("   3 overnight runs lost before sentinel. 5-10 LLM calls/run saves hours.", FG, False),
        ("8. Templates beat conventions", GREEN, True),
        ("   \"Make a Next.js project\" = 5 structures. Templates = 0% divergence.", FG, False),
    ])

    table_slide(prs, "Scaling: MiniShop vs CraftBrew",
        ["Metric", "MiniShop", "CraftBrew #7", "Multiplier"],
        [
            ["Changes", "6", "15", "2.5x"],
            ["Source files", "47", "150+", "3x"],
            ["DB models", "~8", "28", "3.5x"],
            ["Merge conflicts", "0", "4 (all resolved)", "—"],
            ["Human intervention", "0", "0", "—"],
            ["Wall time", "1h 45m", "~6h", "3.4x"],
            ["Tokens", "2.7M", "~11M", "4x"],
        ], col_widths=[3, 2.5, 3, 2.5],
        note="Token scaling is super-linear (4x tokens for 2.5x changes) — later changes need more context.")

    table_slide(prs, "Convergence — measuring reproducibility",
        ["Dimension", "Match"],
        [
            ["DB schema (models, fields, relations)", "100%"],
            ["Conventions (naming, structure)", "100%"],
            ["Routes (URLs, API endpoints)", "83%"],
            ["Overall", "83/100"],
        ], col_widths=[6, 5],
        note="Remaining 17% is stylistic, not structural. Schema and conventions are fully deterministic.")

    # ── Roadmap ──
    section_slide(prs, "Roadmap", "Where we're heading")

    content_slide(prs, "Development priorities", [
        ("Divergence reduction — template optimization, scaffold testing", FG, False),
        ("Build time optimization — parallel gates, incremental builds", FG, False),
        ("Session context reuse — reduce cold-start token overhead", FG, False),
        ("Memory optimization — relevance scoring, auto-rule conversion", FG, False),
        ("Gate intelligence — adaptive thresholds from historical data", FG, False),
        ("Merge conflict prevention — proactive file-level detection", FG, False),
        ("", FG, False),
        ("Core/Web separation — 170+ web refs leaked into core", AMBER, True),
        ("NIS2 Compliance Layer — EU 2022/2555 directive", AMBER, True),
        ("shadcn/ui Design Connector — local design-system.md", AMBER, True),
    ])

    # ── Summary ──
    title_slide(prs, "Summary",
                "SPECIFY → DECOMPOSE → EXECUTE → SUPERVISE → VERIFY → LEARN",
                "6 changes | 1h 45m | 0 interventions | 70 tests | 100% spec coverage\nGreenfield, brownfield, isolated unit — same pipeline.")

    # ── Questions ──
    s = content_slide(prs, "Questions?", [
        ("GitHub: github.com/tatargabor/set-core", GREEN, False),
        ("Website: setcode.dev", GREEN, False),
        ("Benchmarks: docs/learn/benchmarks.md", MUTED, False),
        ("Lessons learned: docs/learn/lessons-learned.md", MUTED, False),
        ("", FG, False),
        ("Try it:", AMBER, True),
        ("pip install -e . && pip install -e modules/web", FG, False),
        ("./tests/e2e/runners/run-micro-web.sh", FG, False),
        ("open http://localhost:7400", FG, False),
    ])

    prs.save(str(output_path))
    print(f"Saved: {output_path} ({output_path.stat().st_size / 1024 / 1024:.1f} MB)")


def main():
    parser = argparse.ArgumentParser(description="Build SET presentation PPTX")
    parser.add_argument("--lang", default="en", choices=["en"])
    parser.add_argument("--output", default=None)
    args = parser.parse_args()

    if args.output:
        out = Path(args.output)
    else:
        out = PRES_DIR / "set-core-presentation-pro.pptx"

    build_en(out)


if __name__ == "__main__":
    main()
